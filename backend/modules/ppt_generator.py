"""
PowerPoint Report Generator â€” IntelliSecure
============================================
Generates executive-level PowerPoint presentations using python-pptx.
Ideal for management briefings and security review meetings.

Slides:
  1. Title Slide
  2. Security Score & Risk Level
  3. Active Threat Summary
  4. Top Attack Sources
  5. Priority Recommendations
  6. Closing / Next Steps

Output: reports/ppt/security_report_YYYYMMDD.pptx

Author: IntelliSecure Team
"""

import os
import sys
import datetime

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

ROOT_DIR = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
PPT_DIR  = os.path.join(ROOT_DIR, "reports", "ppt")

try:
    from pptx                  import Presentation
    from pptx.util             import Inches, Pt, Emu
    from pptx.dml.color        import RGBColor
    from pptx.enum.text        import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


# â”€â”€â”€ Color Constants â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
C_DARK  = RGBColor(0x0F, 0x17, 0x2A)   # Dark navy
C_BLUE  = RGBColor(0x1E, 0x40, 0xAF)   # Primary blue
C_CYAN  = RGBColor(0x06, 0xB6, 0xD4)   # Accent cyan
C_RED   = RGBColor(0xEF, 0x44, 0x44)   # Danger red
C_GREEN = RGBColor(0x10, 0xB9, 0x81)   # Success green
C_AMB   = RGBColor(0xF5, 0x9E, 0x0B)   # Amber warning
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_GREY  = RGBColor(0x94, 0xA3, 0xB8)


def _add_text_box(slide, text: str, left, top, width, height,
                  font_size: int = 14, bold: bool = False,
                  color: RGBColor = C_WHITE, align=PP_ALIGN.LEFT):
    """Helper to add a styled text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def _fill_background(slide, color: RGBColor):
    """Set slide background to a solid color."""
    from pptx.oxml.ns import qn
    from lxml import etree
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def generate_security_ppt(stats: dict, alerts: list, score: dict,
                           top_sources: list, recommendations: dict) -> str:
    """
    Generate a PowerPoint presentation summarizing the security posture.

    Args:
        stats:           Dashboard KPI stats
        alerts:          List of active alerts
        score:           Security score dict
        top_sources:     Top attack sources list
        recommendations: Prioritized recommendations dict

    Returns:
        Absolute path to the generated PPTX file.
    """
    if not PPTX_AVAILABLE:
        raise RuntimeError("python-pptx is not installed. Run: pip install python-pptx")

    os.makedirs(PPT_DIR, exist_ok=True)
    timestamp = datetime.datetime.utcnow().strftime("%Y%m%d_%H%M%S")
    filename  = f"security_report_{timestamp}.pptx"
    filepath  = os.path.join(PPT_DIR, filename)

    prs    = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # Blank layout

    # â”€â”€â”€ Slide 1: Title â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    slide = prs.slides.add_slide(blank_layout)
    _fill_background(slide, C_DARK)

    _add_text_box(slide, "ðŸ”’ INTELISECURE",
                  Inches(1), Inches(1.5), Inches(11), Inches(1.2),
                  font_size=40, bold=True, color=C_CYAN, align=PP_ALIGN.CENTER)

    _add_text_box(slide, "AI-Powered Cybersecurity â€” Executive Security Report",
                  Inches(1), Inches(2.8), Inches(11), Inches(0.8),
                  font_size=20, color=C_WHITE, align=PP_ALIGN.CENTER)

    _add_text_box(slide,
                  f"Report Date: {datetime.datetime.utcnow().strftime('%B %d, %Y')}",
                  Inches(1), Inches(4.5), Inches(11), Inches(0.6),
                  font_size=14, color=C_GREY, align=PP_ALIGN.CENTER)

    # â”€â”€â”€ Slide 2: Security Score â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    slide = prs.slides.add_slide(blank_layout)
    _fill_background(slide, C_DARK)

    _add_text_box(slide, "Security Health Score",
                  Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                  font_size=24, bold=True, color=C_CYAN)

    score_val  = score.get("score", 0)
    risk_level = score.get("risk_level", "Unknown")
    risk_color = C_GREEN if score_val >= 80 else C_AMB if score_val >= 60 else C_RED

    _add_text_box(slide, str(score_val),
                  Inches(4.5), Inches(1.5), Inches(4), Inches(2.5),
                  font_size=90, bold=True, color=risk_color, align=PP_ALIGN.CENTER)

    _add_text_box(slide, f"/ 100  â€”  Risk Level: {risk_level}",
                  Inches(2.5), Inches(4.2), Inches(8), Inches(0.8),
                  font_size=18, color=risk_color, align=PP_ALIGN.CENTER)

    # KPIs row
    kpis = [
        (str(stats.get("total_logs", 0)),      "Total Logs"),
        (str(stats.get("active_threats", 0)),  "Active Threats"),
        (str(stats.get("critical_alerts", 0)), "Critical Alerts"),
    ]
    for i, (val, label) in enumerate(kpis):
        x = Inches(1 + i * 3.7)
        _add_text_box(slide, val,   x, Inches(5.2), Inches(3), Inches(0.8),
                      font_size=28, bold=True, color=C_WHITE)
        _add_text_box(slide, label, x, Inches(6.0), Inches(3), Inches(0.5),
                      font_size=12, color=C_GREY)

    # â”€â”€â”€ Slide 3: Active Threats â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    slide = prs.slides.add_slide(blank_layout)
    _fill_background(slide, C_DARK)

    _add_text_box(slide, "Active Security Threats",
                  Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                  font_size=24, bold=True, color=C_CYAN)

    if alerts:
        for i, alert in enumerate(alerts[:6]):
            sev   = alert.get("severity", "Medium")
            color = C_RED if sev in ("Critical", "High") else C_AMB
            y     = Inches(1.2 + i * 0.95)
            _add_text_box(slide,
                          f"[{sev}]  {alert.get('threat_type', 'â€”')}  â†  {alert.get('source', 'â€”')}",
                          Inches(0.5), y, Inches(12), Inches(0.7),
                          font_size=13, color=color)
    else:
        _add_text_box(slide, "âœ…  No active threats at time of report.",
                      Inches(0.5), Inches(3), Inches(12), Inches(1),
                      font_size=18, color=C_GREEN, align=PP_ALIGN.CENTER)

    # â”€â”€â”€ Slide 4: Top Attack Sources â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    slide = prs.slides.add_slide(blank_layout)
    _fill_background(slide, C_DARK)

    _add_text_box(slide, "Top Attack Sources",
                  Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                  font_size=24, bold=True, color=C_CYAN)

    if top_sources:
        for i, src in enumerate(top_sources[:5]):
            y = Inches(1.3 + i * 1.0)
            _add_text_box(slide, f"#{i+1}  {src.get('source', 'â€”')}",
                          Inches(0.5), y, Inches(7), Inches(0.6),
                          font_size=16, bold=True, color=C_WHITE)
            _add_text_box(slide, f"{src.get('alert_count', 0)} alerts",
                          Inches(8), y, Inches(4), Inches(0.6),
                          font_size=14, color=C_RED)
    else:
        _add_text_box(slide, "No attacker data available.",
                      Inches(0.5), Inches(3), Inches(12), Inches(1),
                      font_size=18, color=C_GREY, align=PP_ALIGN.CENTER)

    # â”€â”€â”€ Slide 5: Recommendations â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    slide = prs.slides.add_slide(blank_layout)
    _fill_background(slide, C_DARK)

    _add_text_box(slide, "Priority Recommendations",
                  Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                  font_size=24, bold=True, color=C_CYAN)

    all_p1_p2 = recommendations.get("p1_critical", []) + recommendations.get("p2_high", [])
    for i, rec in enumerate(all_p1_p2[:6]):
        y = Inches(1.2 + i * 0.95)
        _add_text_box(slide, f"[{rec.get('priority', 'P2')}]  {rec.get('action', '')}",
                      Inches(0.5), y, Inches(12), Inches(0.7),
                      font_size=13, color=C_WHITE)

    # â”€â”€â”€ Slide 6: Next Steps â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    slide = prs.slides.add_slide(blank_layout)
    _fill_background(slide, C_BLUE)

    _add_text_box(slide, "Next Steps & Closing",
                  Inches(0.5), Inches(0.8), Inches(12), Inches(0.8),
                  font_size=28, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    next_steps = [
        "1. Review and act on all P1 Critical recommendations immediately.",
        "2. Schedule follow-up scan within 24 hours.",
        "3. Brief the security team on active threat landscape.",
        "4. Update incident response playbooks if new attack patterns emerged.",
        "5. Re-run IntelliSecure security scan after remediation."
    ]
    for i, step in enumerate(next_steps):
        _add_text_box(slide, step,
                      Inches(1), Inches(1.8 + i * 0.85), Inches(11), Inches(0.7),
                      font_size=14, color=C_WHITE)

    _add_text_box(slide, f"IntelliSecure v2.0 â€” Confidential | {datetime.datetime.utcnow().strftime('%Y-%m-%d')}",
                  Inches(0.5), Inches(6.8), Inches(12), Inches(0.5),
                  font_size=10, color=RGBColor(0xBF, 0xDB, 0xFF), align=PP_ALIGN.CENTER)

    prs.save(filepath)
    return filepath
